#!/usr/bin/env python3
"""
6. スライド生成（PowerPoint版）
slides_plan.jsonをもとにPowerPointテンプレートを使ってスライドを生成
テンプレートファイルをベースにして、スライドを複製・編集する方式
"""

import json
import sys
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
import shutil
from lxml import etree

def load_json(filepath):
    """JSONファイルを読み込み"""
    with open(filepath, 'r', encoding='utf-8') as f:
        return json.load(f)

def get_template_slide_index(template_name, item_count=0):
    """
    テンプレート名からテンプレートスライドのインデックスを取得
    slide_templates_all_variations_jp.pptxの構造:
    スライド1 (index 0): strong-title (強調メッセージ)
    スライド2 (index 1): タイトル+テキスト3行
    スライド3 (index 2): タイトル+テキスト4行
    スライド4 (index 3): タイトル+テキスト5行
    スライド5 (index 4): タイトル+大きなイラスト枠（説明なし）
    スライド6 (index 5): タイトル+大きなイラスト枠+説明1
    スライド7 (index 6): タイトル+イラスト枠2枚（説明なし）
    スライド8 (index 7): タイトル+イラスト枠2枚+説明1,2
    スライド9 (index 8): タイトル+大きなスクリーンショット枠
    スライド10 (index 9): タイトル+大きなスクリーンショット枠+説明1
    """
    template_map = {
        'title_card': 0,      # スライド1: 強調メッセージ
        'cta': 0,             # スライド1: 強調メッセージ
        'definition': 1,      # スライド2: タイトル+テキスト3行
        'bullets': None,      # item数に応じて選択
        'process': None,      # step数に応じて選択
        'recap': None,        # point数に応じて選択
        'comparison': 1,      # スライド2: タイトル+テキスト3行
        'illustration': 4,    # スライド5: タイトル+イラスト枠
        'diagram': 4,         # スライド5: タイトル+イラスト枠
        'screenshot': 8,      # スライド9: タイトル+スクリーンショット枠
    }

    base_idx = template_map.get(template_name)

    # bullets/process/recapは項目数に応じて選択
    if base_idx is None:
        if item_count <= 3:
            return 1  # スライド2: 3行
        elif item_count == 4:
            return 2  # スライド3: 4行
        else:
            return 3  # スライド4: 5行

    return base_idx

def duplicate_slide(prs, slide_index):
    """
    指定したインデックスのスライドを完全に複製（背景を含む）
    背景をXMLレベルでコピーし、図形は通常の方法でコピー
    """
    source_slide = prs.slides[slide_index]

    # 同じレイアウトを使用
    slide_layout = source_slide.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)

    # XML要素を取得
    source_slide_element = source_slide.element
    new_slide_element = new_slide.element

    # 名前空間
    ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}

    # 背景要素をコピー（cSldの中のbg要素）
    source_cSld = source_slide_element.find('.//p:cSld', ns)
    new_cSld = new_slide_element.find('.//p:cSld', ns)

    if source_cSld is not None and new_cSld is not None:
        # 背景要素を探す
        source_bg = source_cSld.find('./p:bg', ns)

        if source_bg is not None:
            # 既存の背景要素を削除
            existing_bg = new_cSld.find('./p:bg', ns)
            if existing_bg is not None:
                new_cSld.remove(existing_bg)

            # 背景要素をコピーして挿入
            new_bg = etree.fromstring(etree.tostring(source_bg))
            # spTreeの前に挿入
            spTree = new_cSld.find('./p:spTree', ns)
            if spTree is not None:
                spTree_index = list(new_cSld).index(spTree)
                new_cSld.insert(spTree_index, new_bg)

    # 図形を複製
    for shape in source_slide.shapes:
        el = shape.element
        newel = etree.fromstring(etree.tostring(el))
        new_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

    return new_slide

def ensure_white_text(run_element):
    """runのテキストを白色に設定（XMLレベル）"""
    from lxml import etree
    ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}

    # rPrを取得または作成
    rPr = run_element.find('.//a:rPr', ns)
    if rPr is None:
        rPr = etree.Element('{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
        run_element.insert(0, rPr)

    # 既存のsolidFillを削除
    existing_solidFill = rPr.find('.//a:solidFill', ns)
    if existing_solidFill is not None:
        rPr.remove(existing_solidFill)

    # 白色のsolidFillを追加
    solidFill = etree.SubElement(rPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
    srgbClr = etree.SubElement(solidFill, '{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
    srgbClr.set('val', 'FFFFFF')

def set_shape_text(shape, text):
    """図形のテキストを設定（単一テキスト用）- フォント書式を保持、白色を強制"""
    try:
        if hasattr(shape, 'text_frame'):
            text_frame = shape.text_frame

            # 最初のパラグラフがあるか確認
            if len(text_frame.paragraphs) > 0:
                para = text_frame.paragraphs[0]

                # 既存のrunがあれば、最初のrunのフォント書式を保持
                if len(para.runs) > 0:
                    # 最初のrunのフォント書式を保存
                    first_run = para.runs[0]
                    font_props = {
                        'name': first_run.font.name,
                        'size': first_run.font.size,
                        'bold': first_run.font.bold,
                        'color_type': first_run.font.color.type if hasattr(first_run.font, 'color') else None
                    }

                    # 色情報を保存（XMLレベル）
                    if hasattr(first_run.font, 'color'):
                        try:
                            if first_run.font.color.type == 1:  # RGB
                                font_props['rgb'] = first_run.font.color.rgb
                            elif first_run.font.color.type == 2:  # SCHEME
                                font_props['theme_color'] = first_run.font.color.theme_color
                        except:
                            pass

                    # runのXML要素を保存
                    run_element = first_run._r

                    # パラグラフのテキストをクリア
                    for run in para.runs:
                        run.text = ""
                    para.text = ""

                    # 新しいテキストを設定
                    para.text = str(text)

                    # フォント書式を再適用
                    if len(para.runs) > 0:
                        new_run = para.runs[0]
                        new_run_element = new_run._r
                        from lxml import etree
                        ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}

                        # XML要素のrPr（run properties）をコピー
                        if run_element is not None:
                            source_rPr = run_element.find('.//a:rPr', ns)

                            # 既存のrPrを削除
                            existing_rPr = new_run_element.find('.//a:rPr', ns)
                            if existing_rPr is not None:
                                existing_rPr.getparent().remove(existing_rPr)

                            if source_rPr is not None:
                                # テンプレートにrPrがある場合はコピー
                                new_rPr = etree.fromstring(etree.tostring(source_rPr))
                                new_run_element.insert(0, new_rPr)
                            # テンプレートにrPrがない場合は、runレベルのプロパティを持たない
                            # （パラグラフのdefRPrから継承される）
                        else:
                            # run要素がない場合も、runレベルのrPrを削除
                            existing_rPr = new_run_element.find('.//a:rPr', ns)
                            if existing_rPr is not None:
                                existing_rPr.getparent().remove(existing_rPr)

                        # 白色を強制適用
                        ensure_white_text(new_run_element)
                else:
                    # runがない場合は通常設定
                    para.text = str(text)
                    # 白色を適用
                    if len(para.runs) > 0:
                        ensure_white_text(para.runs[0]._r)
            else:
                # パラグラフがない場合は追加
                p = text_frame.add_paragraph()
                p.text = str(text)
                # 白色を適用
                if len(p.runs) > 0:
                    ensure_white_text(p.runs[0]._r)
    except Exception as e:
        print(f"Warning: Could not set text for shape: {e}")
        import traceback
        traceback.print_exc()

def set_shape_text_lines(shape, lines):
    """図形のテキストを複数行で設定（各行を別パラグラフに）- フォント書式を保持、白色を強制"""
    try:
        if hasattr(shape, 'text_frame'):
            text_frame = shape.text_frame
            from lxml import etree

            # 各パラグラフのフォント書式を保存
            para_formats = []
            for para in text_frame.paragraphs:
                if len(para.runs) > 0:
                    first_run = para.runs[0]
                    # runのXML要素を保存
                    run_element = etree.fromstring(etree.tostring(first_run._r))
                    para_formats.append(run_element)
                else:
                    para_formats.append(None)

            # 既存のパラグラフをクリア
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.text = ""
                paragraph.text = ""

            # 既存のパラグラフを削除（最初のものは残す）
            while len(text_frame.paragraphs) > 1:
                elem = text_frame.paragraphs[1]._element
                elem.getparent().remove(elem)

            # 各行を個別のパラグラフとして追加
            for idx, line in enumerate(lines):
                if idx == 0:
                    # 最初のパラグラフを使用
                    para = text_frame.paragraphs[0]
                    para.text = str(line)
                else:
                    # 新しいパラグラフを追加
                    para = text_frame.add_paragraph()
                    para.text = str(line)

                # フォント書式を再適用（保存されていれば）
                if len(para.runs) > 0:
                    new_run = para.runs[0]
                    new_run_element = new_run._r
                    ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}

                    # 使用するフォーマットを決定
                    source_run_element = None
                    if idx < len(para_formats) and para_formats[idx] is not None:
                        source_run_element = para_formats[idx]
                    elif len(para_formats) > 0 and para_formats[0] is not None:
                        source_run_element = para_formats[0]

                    if source_run_element is not None:
                        source_rPr = source_run_element.find('.//a:rPr', ns)

                        # 既存のrPrを削除
                        existing_rPr = new_run_element.find('.//a:rPr', ns)
                        if existing_rPr is not None:
                            existing_rPr.getparent().remove(existing_rPr)

                        if source_rPr is not None:
                            # テンプレートにrPrがある場合はコピー
                            new_rPr = etree.fromstring(etree.tostring(source_rPr))
                            new_run_element.insert(0, new_rPr)
                        # テンプレートにrPrがない場合は、runレベルのプロパティを持たない
                    else:
                        # フォーマットがない場合も、runレベルのrPrを削除
                        existing_rPr = new_run_element.find('.//a:rPr', ns)
                        if existing_rPr is not None:
                            existing_rPr.getparent().remove(existing_rPr)

                    # 白色を強制適用
                    ensure_white_text(new_run_element)
    except Exception as e:
        print(f"Warning: Could not set text lines for shape: {e}")
        import traceback
        traceback.print_exc()

def fill_slide_content(slide, fields, template_idx):
    """スライドの内容を埋める"""
    shapes = list(slide.shapes)

    if template_idx == 0:
        # 強調メッセージスライド (shapes[0] = メッセージ)
        if len(shapes) > 0:
            message = fields.get('title') or fields.get('message', '')
            if fields.get('subtitle'):
                message = f"{fields['title']}\n{fields['subtitle']}"
            set_shape_text(shapes[0], message)

    elif template_idx in [1, 2, 3]:
        # リストスライド（shapes[0] = タイトル, shapes[1] = コンテンツ）
        if len(shapes) >= 2:
            # タイトル
            title = fields.get('title', '')
            if fields.get('term'):
                title = fields['term']
            set_shape_text(shapes[0], title)

            # コンテンツ（各行を個別のパラグラフに）
            max_items = [3, 4, 5][template_idx - 1]
            content_lines = []

            if 'items' in fields:
                content_lines = fields['items'][:max_items]
            elif 'steps' in fields:
                content_lines = [f"{i+1}. {step}" for i, step in enumerate(fields['steps'][:max_items])]
            elif 'points' in fields:
                content_lines = fields['points'][:max_items]
            elif 'desc' in fields:
                content_lines = [fields['desc']]

            # 複数行を個別のパラグラフとして設定
            set_shape_text_lines(shapes[1], content_lines)

    elif template_idx in [4, 5, 6, 7, 8, 9]:
        # イラスト/スクリーンショットスライド
        if len(shapes) >= 1:
            title = fields.get('title', '')
            set_shape_text(shapes[0], title)

def generate_pptx(slides_plan_path, template_path, output_path):
    """PowerPointスライドを生成"""
    # slides_plan.jsonまたはtuned.jsonを読み込み
    plan_data = load_json(slides_plan_path)

    # チューニング済みの場合は slidesWithTuning を使用
    if 'slidesWithTuning' in plan_data:
        slides_data = plan_data['slidesWithTuning']
    else:
        slides_data = plan_data.get('slides', [])

    # テンプレートを読み込み
    if not os.path.exists(template_path):
        print(f"Error: Template file not found: {template_path}")
        sys.exit(1)

    # テンプレートをコピーして作業用ファイルを作成
    temp_path = output_path + '.temp.pptx'
    shutil.copy2(template_path, temp_path)

    try:
        # 作業用ファイルを開く
        prs = Presentation(temp_path)

        # テンプレートスライドの数を保存
        num_template_slides = len(prs.slides)
        print(f"Template has {num_template_slides} slides")

        # 各スライドプランに対してスライドを生成
        new_slides = []
        for idx, slide_plan in enumerate(slides_data):
            template_name = slide_plan.get('template', 'bullets')
            fields = slide_plan.get('fields', {})

            # 項目数を取得
            item_count = 0
            if 'items' in fields:
                item_count = len(fields['items'])
            elif 'steps' in fields:
                item_count = len(fields['steps'])
            elif 'points' in fields:
                item_count = len(fields['points'])

            # テンプレートスライドのインデックスを取得
            template_idx = get_template_slide_index(template_name, item_count)
            print(f"Slide {idx + 1}: Using template {template_idx + 1} for '{template_name}' with {item_count} items")

            # テンプレートスライドを複製
            if template_idx < num_template_slides:
                new_slide = duplicate_slide(prs, template_idx)
                new_slides.append(new_slide)

                # 内容を埋める
                fill_slide_content(new_slide, fields, template_idx)
            else:
                print(f"Warning: Template index {template_idx} out of range")

        # 元のテンプレートスライドを削除
        print(f"Removing {num_template_slides} template slides...")
        for i in range(num_template_slides):
            # 常に最初のスライドを削除（削除するとインデックスがずれるため）
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]

        # 出力ディレクトリを作成
        output_dir = os.path.dirname(output_path)
        if output_dir:
            os.makedirs(output_dir, exist_ok=True)

        # PowerPointファイルを保存
        prs.save(output_path)
        print(f"Generated PowerPoint: {output_path} ({len(slides_data)} slides)")

    except Exception as e:
        print(f"Error generating PowerPoint: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)

    finally:
        # 一時ファイルを削除
        if os.path.exists(temp_path):
            os.remove(temp_path)

def main():
    if len(sys.argv) < 4:
        print("Usage: python src/06_render_pptx.py <slides_plan.json> <template.pptx> <output.pptx>")
        sys.exit(1)

    slides_plan_path = sys.argv[1]
    template_path = sys.argv[2]
    output_path = sys.argv[3]

    generate_pptx(slides_plan_path, template_path, output_path)

if __name__ == '__main__':
    main()
