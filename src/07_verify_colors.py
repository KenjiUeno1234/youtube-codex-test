#!/usr/bin/env python3
"""
7. 文字色検証と修正スクリプト
生成されたPowerPointファイルの全テキストが白色になっているか検証し、
必要に応じて白色に修正する
"""

import sys
import os
from pptx import Presentation
from lxml import etree

def verify_and_fix_text_colors(pptx_path, output_path=None):
    """
    PowerPointファイルの全テキスト色を検証し、白色でない場合は修正する

    Args:
        pptx_path: 検証対象のPowerPointファイルパス
        output_path: 出力先（Noneの場合は上書き）

    Returns:
        dict: 検証結果の統計情報
    """
    if not os.path.exists(pptx_path):
        print(f"Error: File not found: {pptx_path}")
        sys.exit(1)

    if output_path is None:
        output_path = pptx_path

    prs = Presentation(pptx_path)

    stats = {
        'total_slides': len(prs.slides),
        'total_shapes': 0,
        'total_paragraphs': 0,
        'total_runs': 0,
        'fixed_runs': 0,
        'already_white': 0,
        'no_color': 0,
        'issues': []
    }

    ns = {
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'
    }

    print(f"\n=== Color Verification Start: {pptx_path} ===\n")

    for slide_idx, slide in enumerate(prs.slides, 1):
        print(f"Slide {slide_idx}:")

        for shape_idx, shape in enumerate(slide.shapes):
            stats['total_shapes'] += 1

            if not hasattr(shape, 'text_frame'):
                continue

            text_frame = shape.text_frame

            for para_idx, paragraph in enumerate(text_frame.paragraphs):
                stats['total_paragraphs'] += 1

                for run_idx, run in enumerate(paragraph.runs):
                    stats['total_runs'] += 1
                    run_element = run._r

                    # rPr（run properties）を取得または作成
                    rPr = run_element.find('.//a:rPr', ns)

                    # 現在の色情報を確認
                    current_color = None
                    if rPr is not None:
                        solidFill = rPr.find('.//a:solidFill', ns)
                        if solidFill is not None:
                            # RGBカラー
                            srgbClr = solidFill.find('.//a:srgbClr', ns)
                            if srgbClr is not None:
                                current_color = srgbClr.get('val')

                            # スキームカラー
                            schemeClr = solidFill.find('.//a:schemeClr', ns)
                            if schemeClr is not None:
                                current_color = f"scheme:{schemeClr.get('val')}"

                    # 色の状態を判定
                    if current_color in ['FFFFFF', 'ffffff', 'scheme:lt1', 'scheme:tx1']:
                        stats['already_white'] += 1
                        status = f"OK White ({current_color})"
                        continue
                    else:
                        # 白色以外の色、または色指定なしを修正
                        if current_color is None:
                            stats['no_color'] += 1
                            status = f"FIXED: None -> FFFFFF"
                        else:
                            stats['fixed_runs'] += 1
                            status = f"FIXED: {current_color} -> FFFFFF"

                        # rPrが存在しない場合は作成
                        if rPr is None:
                            rPr = etree.Element('{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
                            run_element.insert(0, rPr)

                        # solidFillを削除して再作成
                        existing_solidFill = rPr.find('.//a:solidFill', ns)
                        if existing_solidFill is not None:
                            rPr.remove(existing_solidFill)

                        # 白色のsolidFillを追加
                        solidFill = etree.SubElement(rPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
                        srgbClr = etree.SubElement(solidFill, '{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
                        srgbClr.set('val', 'FFFFFF')

                        stats['issues'].append({
                            'slide': slide_idx,
                            'shape': shape_idx,
                            'text': run.text[:30] + ('...' if len(run.text) > 30 else ''),
                            'old_color': current_color if current_color else 'None',
                            'new_color': 'FFFFFF'
                        })

                    if run.text.strip():  # テキストがある場合のみ表示
                        text_preview = run.text[:30] + ('...' if len(run.text) > 30 else '')
                        print(f"  Shape {shape_idx}, Para {para_idx}, Run {run_idx}: {status}")
                        print(f"    Text: '{text_preview}'")

    # 結果を保存
    prs.save(output_path)

    # サマリーを表示
    print(f"\n=== Verification Summary ===")
    print(f"Total Slides: {stats['total_slides']}")
    print(f"Total Shapes: {stats['total_shapes']}")
    print(f"Total Paragraphs: {stats['total_paragraphs']}")
    print(f"Total Runs: {stats['total_runs']}")
    print(f"")
    print(f"OK Already white: {stats['already_white']}")
    print(f"FIXED No color -> white: {stats['no_color']}")
    print(f"FIXED Other color -> white: {stats['fixed_runs']}")

    if stats['fixed_runs'] > 0 or stats['no_color'] > 0:
        print(f"\n=== Fix Details ===")
        for issue in stats['issues']:
            print(f"Slide {issue['slide']}, Shape {issue['shape']}: {issue['old_color']} -> {issue['new_color']}")
            print(f"  Text: '{issue['text']}'")

    print(f"\nOutput file: {output_path}")

    return stats

def main():
    if len(sys.argv) < 2:
        print("Usage: python src/07_verify_colors.py <input.pptx> [output.pptx]")
        print("  If output.pptx is not specified, the input file will be overwritten")
        sys.exit(1)

    input_path = sys.argv[1]
    output_path = sys.argv[2] if len(sys.argv) >= 3 else None

    stats = verify_and_fix_text_colors(input_path, output_path)

    # 終了コード（修正があった場合は1を返す）
    total_fixed = stats['fixed_runs'] + stats['no_color']
    exit_code = 1 if total_fixed > 0 else 0

    if exit_code == 0:
        print("\nOK All text colors are correct")
    else:
        print(f"\nOK Fixed {total_fixed} locations")

    sys.exit(exit_code)

if __name__ == '__main__':
    main()
