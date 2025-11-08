#!/usr/bin/env python3
"""
テンプレートと生成ファイルのレイアウト/マスター比較
"""

import sys
from pptx import Presentation

def compare_slide_layouts(template_path, generated_path):
    """レイアウトとマスターを比較"""

    template_prs = Presentation(template_path)
    generated_prs = Presentation(generated_path)

    print("="*70)
    print("テンプレートファイル")
    print("="*70)
    print(f"スライドマスター名: {template_prs.slide_master.name if hasattr(template_prs.slide_master, 'name') else 'N/A'}")

    # テンプレートのスライド2を確認
    if len(template_prs.slides) > 1:
        slide = template_prs.slides[1]
        print(f"\nスライド2:")
        print(f"  レイアウト名: {slide.slide_layout.name}")
        print(f"  follow_master_background: {slide.follow_master_background if hasattr(slide, 'follow_master_background') else 'N/A'}")

        # 背景情報
        if hasattr(slide, 'background'):
            bg = slide.background
            print(f"  背景:")
            print(f"    fill.type: {bg.fill.type if hasattr(bg, 'fill') else 'N/A'}")

    print("\n" + "="*70)
    print("生成ファイル")
    print("="*70)
    print(f"スライドマスター名: {generated_prs.slide_master.name if hasattr(generated_prs.slide_master, 'name') else 'N/A'}")

    # 生成ファイルのスライド2を確認
    if len(generated_prs.slides) > 1:
        slide = generated_prs.slides[1]
        print(f"\nスライド2:")
        print(f"  レイアウト名: {slide.slide_layout.name}")
        print(f"  follow_master_background: {slide.follow_master_background if hasattr(slide, 'follow_master_background') else 'N/A'}")

        # 背景情報
        if hasattr(slide, 'background'):
            bg = slide.background
            print(f"  背景:")
            print(f"    fill.type: {bg.fill.type if hasattr(bg, 'fill') else 'N/A'}")

    # レイアウト比較
    print("\n" + "="*70)
    print("レイアウト比較")
    print("="*70)
    print(f"テンプレート レイアウト数: {len(template_prs.slide_layouts)}")
    print(f"生成ファイル レイアウト数: {len(generated_prs.slide_layouts)}")

    # Blank レイアウトを確認
    for idx, layout in enumerate(template_prs.slide_layouts):
        if 'Blank' in layout.name:
            print(f"\nテンプレート Blank レイアウト (index {idx}):")
            print(f"  名前: {layout.name}")
            if hasattr(layout, 'background'):
                print(f"  背景あり: {hasattr(layout.background, 'fill')}")

if __name__ == '__main__':
    if len(sys.argv) < 3:
        print("Usage: python src/compare_layouts.py <template.pptx> <generated.pptx>")
        sys.exit(1)

    compare_slide_layouts(sys.argv[1], sys.argv[2])
