#!/usr/bin/env python3
"""
テンプレートと生成されたPowerPointを比較
"""

import sys
from pptx import Presentation

def compare_slides(template_path, generated_path):
    """2つのPowerPointファイルのスライドを比較"""

    template_prs = Presentation(template_path)
    generated_prs = Presentation(generated_path)

    print("="*70)
    print(f"テンプレート: {template_path}")
    print(f"  スライド数: {len(template_prs.slides)}")

    print(f"\n生成ファイル: {generated_path}")
    print(f"  スライド数: {len(generated_prs.slides)}")
    print("="*70)

    # テンプレートの最初のスライド
    if len(template_prs.slides) > 1:
        t_slide = template_prs.slides[1]  # スライド2（3行リスト）
        print(f"\nテンプレート スライド2:")
        print(f"  図形数: {len(t_slide.shapes)}")
        for idx, shape in enumerate(t_slide.shapes):
            if hasattr(shape, 'text_frame'):
                print(f"  図形{idx}: パラグラフ数={len(shape.text_frame.paragraphs)}")
                for p_idx, para in enumerate(shape.text_frame.paragraphs):
                    print(f"    パラグラフ{p_idx}: '{para.text}'")

    # 生成されたファイルの最初のスライド
    if len(generated_prs.slides) > 1:
        g_slide = generated_prs.slides[1]
        print(f"\n生成ファイル スライド2:")
        print(f"  図形数: {len(g_slide.shapes)}")
        for idx, shape in enumerate(g_slide.shapes):
            if hasattr(shape, 'text_frame'):
                print(f"  図形{idx}: パラグラフ数={len(shape.text_frame.paragraphs)}")
                for p_idx, para in enumerate(shape.text_frame.paragraphs):
                    print(f"    パラグラフ{p_idx}: '{para.text}'")

if __name__ == '__main__':
    if len(sys.argv) < 3:
        print("Usage: python src/compare_pptx.py <template.pptx> <generated.pptx>")
        sys.exit(1)

    compare_slides(sys.argv[1], sys.argv[2])
