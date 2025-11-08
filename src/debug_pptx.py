#!/usr/bin/env python3
"""
PowerPointファイルの詳細デバッグ
"""

import sys
from pptx import Presentation

def debug_pptx(pptx_path):
    """PowerPointファイルの詳細を表示"""
    prs = Presentation(pptx_path)

    print(f"PowerPoint: {pptx_path}")
    print(f"スライド数: {len(prs.slides)}")
    print(f"レイアウト数: {len(prs.slide_layouts)}\n")

    for slide_idx, slide in enumerate(prs.slides):
        print(f"\n{'='*70}")
        print(f"スライド {slide_idx + 1}:")
        print(f"  レイアウト名: {slide.slide_layout.name}")
        print(f"  図形数: {len(slide.shapes)}")

        for shape_idx, shape in enumerate(slide.shapes):
            print(f"\n  図形 {shape_idx + 1}:")
            print(f"    名前: {shape.name}")
            print(f"    タイプ: {shape.shape_type}")
            print(f"    位置: left={shape.left}, top={shape.top}")
            print(f"    サイズ: width={shape.width}, height={shape.height}")

            # テキスト情報
            if hasattr(shape, 'text'):
                print(f"    テキスト: '{shape.text}'")

            if hasattr(shape, 'text_frame'):
                print(f"    text_frame: あり")
                print(f"    パラグラフ数: {len(shape.text_frame.paragraphs)}")
                for para_idx, para in enumerate(shape.text_frame.paragraphs):
                    print(f"      パラグラフ{para_idx + 1}: '{para.text}'")

if __name__ == '__main__':
    if len(sys.argv) < 2:
        print("Usage: python src/debug_pptx.py <file.pptx>")
        sys.exit(1)

    debug_pptx(sys.argv[1])
