#!/usr/bin/env python3
"""
PowerPointテンプレートの構造を分析
"""

import sys
from pptx import Presentation

def analyze_template(template_path):
    """テンプレートのレイアウトとプレースホルダーを分析"""
    prs = Presentation(template_path)

    print(f"テンプレート分析: {template_path}")
    print(f"総スライド数: {len(prs.slides)}")
    print(f"総レイアウト数: {len(prs.slide_layouts)}\n")

    # 各レイアウトを分析
    for idx, layout in enumerate(prs.slide_layouts):
        print(f"=" * 60)
        print(f"レイアウト {idx}: {layout.name}")
        print(f"プレースホルダー数: {len(layout.placeholders)}")

        for ph_idx, placeholder in enumerate(layout.placeholders):
            print(f"  [{ph_idx}] {placeholder.name}")
            print(f"      Type: {placeholder.placeholder_format.type}")
            print(f"      idx: {placeholder.placeholder_format.idx}")
        print()

    # 既存のスライドを分析
    print(f"\n" + "=" * 60)
    print("既存スライドの分析:")
    for slide_idx, slide in enumerate(prs.slides):
        print(f"\nスライド {slide_idx + 1}:")
        print(f"  レイアウト: {slide.slide_layout.name}")
        print(f"  図形数: {len(slide.shapes)}")

        for shape_idx, shape in enumerate(slide.shapes):
            print(f"    [{shape_idx}] {shape.name}")
            if hasattr(shape, "text"):
                preview = shape.text[:50].replace('\n', ' ') if shape.text else "(空)"
                print(f"        テキスト: {preview}")

if __name__ == '__main__':
    if len(sys.argv) < 2:
        print("Usage: python src/analyze_pptx_template.py <template.pptx>")
        sys.exit(1)

    analyze_template(sys.argv[1])
