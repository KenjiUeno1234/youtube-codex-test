#!/usr/bin/env python3
"""
XMLレベルでパラグラフとrunの完全な構造を分析
"""
import sys
from pptx import Presentation
from lxml import etree

def analyze_shape_xml(pptx_path, slide_idx, shape_idx):
    """指定された図形のXML構造を完全に表示"""
    prs = Presentation(pptx_path)

    if slide_idx >= len(prs.slides):
        print(f"スライド {slide_idx + 1} が見つかりません")
        return

    slide = prs.slides[slide_idx]
    shapes = list(slide.shapes)

    if shape_idx >= len(shapes):
        print(f"図形 {shape_idx} が見つかりません")
        return

    shape = shapes[shape_idx]

    print(f"{'='*70}")
    print(f"ファイル: {pptx_path}")
    print(f"スライド {slide_idx + 1}, 図形 {shape_idx}: {shape.name}")
    print(f"{'='*70}\n")

    if hasattr(shape, 'text_frame'):
        # 図形全体のXMLを表示
        shape_element = shape.element
        xml_str = etree.tostring(shape_element, encoding='unicode', pretty_print=True)
        print("完全なXML:")
        print(xml_str)
    else:
        print("この図形にはtext_frameがありません")

if __name__ == '__main__':
    if len(sys.argv) < 4:
        print("Usage: python src/analyze_xml_full.py <pptx_file> <slide_index> <shape_index>")
        sys.exit(1)

    pptx_path = sys.argv[1]
    slide_idx = int(sys.argv[2])
    shape_idx = int(sys.argv[3])

    analyze_shape_xml(pptx_path, slide_idx, shape_idx)
