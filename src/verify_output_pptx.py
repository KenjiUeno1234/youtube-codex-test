#!/usr/bin/env python3
"""
生成されたPowerPointの内容を検証
"""

import sys
from pptx import Presentation

def verify_output(pptx_path):
    """出力されたPowerPointを検証"""
    prs = Presentation(pptx_path)

    print(f"生成されたPowerPoint: {pptx_path}")
    print(f"総スライド数: {len(prs.slides)}\n")

    for slide_idx, slide in enumerate(prs.slides):
        print(f"=" * 60)
        print(f"スライド {slide_idx + 1}:")
        print(f"  レイアウト: {slide.slide_layout.name}")
        print(f"  図形数: {len(slide.shapes)}")

        for shape_idx, shape in enumerate(slide.shapes):
            print(f"    [{shape_idx}] {shape.name}")
            if hasattr(shape, "text") and shape.text:
                # テキストをプレビュー（最初の100文字）
                preview = shape.text[:100].replace('\n', ' | ')
                print(f"        テキスト: {preview}")

if __name__ == '__main__':
    if len(sys.argv) < 2:
        print("Usage: python src/verify_output_pptx.py <output.pptx>")
        sys.exit(1)

    verify_output(sys.argv[1])
