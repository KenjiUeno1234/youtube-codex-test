#!/usr/bin/env python3
"""
テンプレートと生成ファイルのフォント色を比較
"""

import sys
from pptx import Presentation
from pptx.enum.dml import MSO_THEME_COLOR

def analyze_fonts(pptx_path, slide_index=1):
    """スライドのフォント色を分析"""
    prs = Presentation(pptx_path)

    if slide_index >= len(prs.slides):
        print(f"Error: Slide {slide_index} does not exist")
        return

    slide = prs.slides[slide_index]

    print(f"\n{'='*70}")
    print(f"ファイル: {pptx_path}")
    print(f"スライド {slide_index + 1}")
    print(f"{'='*70}")

    for shape_idx, shape in enumerate(slide.shapes):
        if hasattr(shape, 'text_frame'):
            print(f"\n図形 {shape_idx}: {shape.name}")
            text_frame = shape.text_frame

            for para_idx, para in enumerate(text_frame.paragraphs[:3]):
                if para.runs:
                    for run_idx, run in enumerate(para.runs):
                        font = run.font
                        print(f"  パラグラフ{para_idx} Run{run_idx}:")
                        print(f"    text: '{run.text[:20]}'")
                        print(f"    font.name: {font.name}")
                        print(f"    font.size: {font.size}")
                        print(f"    font.bold: {font.bold}")

                        # 色情報
                        if hasattr(font, 'color'):
                            color = font.color
                            print(f"    color.type: {color.type}")

                            if color.type == 1:  # RGB
                                try:
                                    print(f"    color.rgb: {color.rgb}")
                                except:
                                    print(f"    color.rgb: (取得不可)")
                            elif color.type == 2:  # SCHEME
                                try:
                                    print(f"    color.theme_color: {color.theme_color}")
                                except:
                                    print(f"    color.theme_color: (取得不可)")

if __name__ == '__main__':
    if len(sys.argv) < 2:
        print("Usage: python src/analyze_font_colors.py <file.pptx> [slide_index]")
        sys.exit(1)

    pptx_path = sys.argv[1]
    slide_index = int(sys.argv[2]) if len(sys.argv) > 2 else 1

    analyze_fonts(pptx_path, slide_index)
