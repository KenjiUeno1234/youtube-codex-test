#!/usr/bin/env python3
"""
XMLレベルでrun propertiesを詳細に分析
"""
import sys
from pptx import Presentation
from lxml import etree

def analyze_run_xml(pptx_path, slide_idx):
    """指定されたスライドのrun XMLを詳細に分析"""
    prs = Presentation(pptx_path)

    if slide_idx >= len(prs.slides):
        print(f"スライド {slide_idx + 1} が見つかりません")
        return

    slide = prs.slides[slide_idx]
    print(f"{'='*70}")
    print(f"ファイル: {pptx_path}")
    print(f"スライド {slide_idx + 1}")
    print(f"{'='*70}\n")

    for shape_idx, shape in enumerate(slide.shapes):
        if hasattr(shape, 'text_frame'):
            print(f"図形 {shape_idx}: {shape.name}")
            text_frame = shape.text_frame

            for para_idx, para in enumerate(text_frame.paragraphs):
                print(f"  パラグラフ{para_idx}:")

                for run_idx, run in enumerate(para.runs):
                    print(f"    Run{run_idx}: '{run.text}'")

                    # XMLを取得
                    run_element = run._r
                    xml_str = etree.tostring(run_element, encoding='unicode', pretty_print=True)

                    # rPr要素を探す
                    ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                    rPr = run_element.find('.//a:rPr', ns)

                    if rPr is not None:
                        print(f"      rPr存在: あり")
                        rPr_xml = etree.tostring(rPr, encoding='unicode', pretty_print=True)
                        # インデントして表示
                        for line in rPr_xml.split('\n'):
                            if line.strip():
                                print(f"      {line}")
                    else:
                        print(f"      rPr存在: なし")
                    print()
            print()

if __name__ == '__main__':
    if len(sys.argv) < 2:
        print("Usage: python src/analyze_xml_rpr.py <pptx_file> [slide_index]")
        sys.exit(1)

    pptx_path = sys.argv[1]
    slide_idx = int(sys.argv[2]) if len(sys.argv) > 2 else 0

    analyze_run_xml(pptx_path, slide_idx)
