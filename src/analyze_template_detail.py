#!/usr/bin/env python3
"""
テンプレートファイルの詳細な構造分析
背景、フォーマット、スタイル情報を調査
"""

import sys
from pptx import Presentation
from pptx.util import Pt

def analyze_slide_detail(prs, slide_index):
    """スライドの詳細情報を分析"""
    if slide_index >= len(prs.slides):
        print(f"Error: Slide {slide_index} does not exist")
        return

    slide = prs.slides[slide_index]

    print(f"\n{'='*70}")
    print(f"スライド {slide_index + 1} の詳細分析")
    print(f"{'='*70}")

    # スライドレイアウト情報
    print(f"\nレイアウト名: {slide.slide_layout.name}")

    # 背景情報
    print(f"\n背景情報:")
    if hasattr(slide, 'background'):
        print(f"  background属性: あり")
    if hasattr(slide, 'follow_master_background'):
        print(f"  follow_master_background: {slide.follow_master_background}")

    # 各図形の詳細
    print(f"\n図形の詳細:")
    for idx, shape in enumerate(slide.shapes):
        print(f"\n  図形 {idx}: {shape.name}")
        print(f"    タイプ: {shape.shape_type}")

        # フィル情報
        if hasattr(shape, 'fill'):
            fill = shape.fill
            print(f"    フィル:")
            print(f"      type: {fill.type}")
            if hasattr(fill, 'fore_color'):
                try:
                    print(f"      色: {fill.fore_color.rgb if hasattr(fill.fore_color, 'rgb') else 'N/A'}")
                except:
                    print(f"      色: (取得不可)")

        # テキストフレーム情報
        if hasattr(shape, 'text_frame'):
            tf = shape.text_frame
            print(f"    テキストフレーム:")
            print(f"      パラグラフ数: {len(tf.paragraphs)}")

            for p_idx, para in enumerate(tf.paragraphs[:3]):  # 最初の3つのみ
                print(f"      パラグラフ {p_idx}:")
                print(f"        テキスト: '{para.text[:30]}'")
                print(f"        レベル: {para.level}")
                print(f"        alignment: {para.alignment}")

                # フォント情報
                if para.runs:
                    run = para.runs[0]
                    font = run.font
                    print(f"        フォント:")
                    print(f"          name: {font.name}")
                    print(f"          size: {font.size}")
                    print(f"          bold: {font.bold}")
                    print(f"          color: {font.color.rgb if hasattr(font.color, 'rgb') else 'N/A'}")

def analyze_master(prs):
    """スライドマスターを分析"""
    print(f"\n{'='*70}")
    print(f"スライドマスター分析")
    print(f"{'='*70}")

    master = prs.slide_master
    print(f"  名前: {master.name if hasattr(master, 'name') else 'N/A'}")

    # 背景
    if hasattr(master, 'background'):
        bg = master.background
        print(f"  背景: あり")
        if hasattr(bg, 'fill'):
            print(f"    fill type: {bg.fill.type}")

    # 図形
    print(f"  マスターの図形数: {len(master.shapes)}")
    for idx, shape in enumerate(master.shapes[:5]):  # 最初の5つ
        print(f"    図形{idx}: {shape.name}, type={shape.shape_type}")

def main():
    if len(sys.argv) < 2:
        print("Usage: python src/analyze_template_detail.py <template.pptx> [slide_index]")
        sys.exit(1)

    template_path = sys.argv[1]
    slide_index = int(sys.argv[2]) if len(sys.argv) > 2 else 1

    prs = Presentation(template_path)

    print(f"テンプレート: {template_path}")
    print(f"総スライド数: {len(prs.slides)}")
    print(f"総レイアウト数: {len(prs.slide_layouts)}")

    # スライドマスター分析
    analyze_master(prs)

    # 指定されたスライドを分析
    analyze_slide_detail(prs, slide_index)

if __name__ == '__main__':
    main()
