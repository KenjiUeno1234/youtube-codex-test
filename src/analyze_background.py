#!/usr/bin/env python3
"""
スライドの背景設定を詳細に分析
"""
import sys
from pptx import Presentation
from lxml import etree

def analyze_background(pptx_path, slide_idx):
    """指定されたスライドの背景を分析"""
    prs = Presentation(pptx_path)

    if slide_idx >= len(prs.slides):
        print(f"スライド {slide_idx + 1} が見つかりません")
        return

    slide = prs.slides[slide_idx]
    print(f"{'='*70}")
    print(f"ファイル: {pptx_path}")
    print(f"スライド {slide_idx + 1}")
    print(f"{'='*70}\n")

    # follow_master_backgroundを確認
    print(f"follow_master_background: {slide.follow_master_background}")

    # XMLレベルで背景を確認
    slide_element = slide.element
    ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
          'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}

    cSld = slide_element.find('.//p:cSld', ns)
    if cSld is not None:
        bg = cSld.find('./p:bg', ns)
        if bg is not None:
            print("\n背景要素(p:bg)が存在します:")
            bg_xml = etree.tostring(bg, encoding='unicode', pretty_print=True)
            print(bg_xml)
        else:
            print("\n背景要素(p:bg)が存在しません")
    else:
        print("\ncSld要素が見つかりません")

if __name__ == '__main__':
    if len(sys.argv) < 2:
        print("Usage: python src/analyze_background.py <pptx_file> [slide_index]")
        sys.exit(1)

    pptx_path = sys.argv[1]
    slide_idx = int(sys.argv[2]) if len(sys.argv) > 2 else 0

    analyze_background(pptx_path, slide_idx)
